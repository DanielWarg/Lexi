"""
Lexi PowerPoint Skill
=====================
Generates .pptx presentations from text input.
"""

import os
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Dict
from skills.base_skill import BaseSkill, SkillContext, SkillResult


class PowerPointSkill(BaseSkill):
    """
    Skill for creating PowerPoint presentations.
    
    Takes topic, audience, and preferences to generate
    a complete .pptx file.
    """
    
    name = "powerpoint"
    display_name = "PowerPoint-skapare"
    description = "Skapar professionella presentationer (.pptx)"
    triggers = [
        "skapa presentation",
        "gör en powerpoint",
        "skapa slides",
        "presentation om",
        "create presentation",
        "make slides"
    ]
    
    requires_confirmation = False
    
    def __init__(self):
        super().__init__()
        self.output_dir = Path.home() / "Documents" / "Lexi" / "Presentations"
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    async def execute(self, context: SkillContext) -> SkillResult:
        """Generate a PowerPoint presentation."""
        try:
            # Import python-pptx (installed separately)
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
            
            if context.on_status:
                context.on_status("Analyserar ämne...")
            
            # Parse input to extract topic and structure
            topic = self._extract_topic(context.user_input)
            num_slides = self._extract_slide_count(context.user_input)
            
            if context.on_status:
                context.on_status(f"Skapar presentation om '{topic}'...")
            
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # Generate slides using LLM context
            slides_content = self._generate_slides_content(topic, num_slides, context)
            
            if context.on_progress:
                context.on_progress(20, "Lägger till slides...")
            
            for i, slide_data in enumerate(slides_content):
                self._add_slide(prs, slide_data, i == 0)
                
                if context.on_progress:
                    progress = 20 + int((i + 1) / len(slides_content) * 70)
                    context.on_progress(progress, f"Slide {i + 1}/{len(slides_content)}")
            
            # Save file
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_topic = "".join(c for c in topic if c.isalnum() or c in " -_")[:30]
            filename = f"{safe_topic}_{timestamp}.pptx"
            filepath = self.output_dir / filename
            
            if context.on_status:
                context.on_status("Sparar presentation...")
            
            prs.save(str(filepath))
            
            if context.on_progress:
                context.on_progress(100, "Klart!")
            
            return SkillResult(
                success=True,
                data={
                    "slides": len(slides_content),
                    "topic": topic
                },
                file_path=str(filepath),
                message=f"Presentation skapad: {filename}"
            )
            
        except ImportError:
            return SkillResult(
                success=False,
                error="python-pptx är inte installerat. Kör: pip install python-pptx"
            )
        except Exception as e:
            return SkillResult(
                success=False,
                error=f"PowerPoint-fel: {str(e)}"
            )
    
    def _extract_topic(self, text: str) -> str:
        """Extract the presentation topic from user input."""
        # Remove trigger words
        triggers = ["skapa presentation", "gör en powerpoint", "presentation om", "slides om"]
        result = text.lower()
        for trigger in triggers:
            result = result.replace(trigger, "")
        return result.strip() or "Presentation"
    
    def _extract_slide_count(self, text: str) -> int:
        """Extract desired number of slides from input."""
        import re
        match = re.search(r'(\d+)\s*slides?', text.lower())
        if match:
            return min(int(match.group(1)), 20)  # Max 20 slides
        return 6  # Default
    
    def _generate_slides_content(
        self, 
        topic: str, 
        num_slides: int,
        context: SkillContext
    ) -> List[Dict]:
        """
        Generate slide content.
        In production, this would call the LLM for content generation.
        For now, returns a template structure.
        """
        slides = [
            {"title": topic.title(), "content": ["Presentation skapad av Lexi"], "is_title": True}
        ]
        
        # Add content slides
        sections = ["Bakgrund", "Huvudpunkter", "Analys", "Slutsats"]
        for i, section in enumerate(sections[:num_slides - 2]):
            slides.append({
                "title": section,
                "content": [
                    f"Punkt 1 om {topic}",
                    f"Punkt 2 om {topic}",
                    f"Punkt 3 om {topic}"
                ],
                "is_title": False
            })
        
        # Add closing slide
        slides.append({
            "title": "Tack!",
            "content": ["Frågor?"],
            "is_title": False
        })
        
        return slides
    
    def _add_slide(self, prs, slide_data: Dict, is_title_slide: bool):
        """Add a slide to the presentation."""
        from pptx.util import Inches, Pt
        
        if is_title_slide:
            layout = prs.slide_layouts[6]  # Blank
        else:
            layout = prs.slide_layouts[6]  # Blank
        
        slide = prs.slides.add_slide(layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data["title"]
        title_para.font.size = Pt(44 if is_title_slide else 32)
        title_para.font.bold = True
        
        # Add content
        if slide_data.get("content"):
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2), Inches(12), Inches(5)
            )
            content_frame = content_box.text_frame
            
            for i, point in enumerate(slide_data["content"]):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()
                para.text = f"• {point}"
                para.font.size = Pt(24 if is_title_slide else 20)
                para.space_after = Pt(12)
